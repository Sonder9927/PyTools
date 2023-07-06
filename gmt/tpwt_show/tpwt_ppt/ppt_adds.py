from pptx import util
from pathlib import Path
import json


def ppt_add_diffs(prs, figs: Path, info_file, margin):
    """
    This script can insert a diff picture
    with info into a new slide of pptx.
    """
    # read vel info json
    with open(info_file) as f:
        info = json.load(f)

    [left, top] = margin
    # [width, height] = shape

    # insert title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Diff Results"

    blank_slide_layout = prs.slide_layouts[6]
    # insert dispersion curves
    dc = figs / r"dispersion_curves.png"
    if dc.exists():
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(str(dc), util.Cm(2.5), util.Cm(5),
                                 util.Cm(20), util.Cm(10))

    for f in figs.glob("*diff*"):
        # insert diff fig into new slide
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(str(f), util.Cm(left), util.Cm(top))
        # insert table
        per = f.stem.split("_")[-1]
        table_rows, table_cols = 3, 2
        table_left, table_top = util.Inches(4.5), util.Inches(4)
        table_width, table_height = util.Inches(3.5), util.Inches(0.8)
        table = slide.shapes.add_table(table_rows, table_cols, table_left,
                                       table_top, table_width,
                                       table_height).table

        table.style = "VALUE TABLE"
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"

        table.cell(1, 0).text = "mean"
        table.cell(1, 1).text = info[per]["avg_diff"]
        table.cell(2, 0).text = "std"
        table.cell(2, 1).text = info[per]["standard_deviation"]

    return prs


def ppt_add_single_type(prs, figs, margin, shape, type):
    """
    This script can insert r*c pictures
    into a new slide of pptx.
    """
    nrow, ncol = 2, 3  # shape
    irow, icol = .382, 1  # itvel
    [left_margin, top_margin] = margin
    [width, height] = shape

    # insert title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = f"{type[1:-1]} Results"

    # make batch figs for every slide
    batch = nrow * ncol
    vel_figs = sorted(figs.glob(type),
                      key=lambda p: int(p.stem.split("_")[-1]))
    figs_batch = [
        vel_figs[i:i + batch] for i in range(0, len(vel_figs), batch)
    ]
    blank_slide_layout = prs.slide_layouts[6]
    for fb in figs_batch:
        # new slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # add pictures
        for i, f in enumerate(fb):
            # set margin for each fig
            left = util.Cm(left_margin + i % ncol * (width + icol))
            top = util.Cm(top_margin + i // ncol * (height + irow))

            slide.shapes.add_picture(str(f), left, top, util.Cm(width),
                                     util.Cm(height))

    return prs
