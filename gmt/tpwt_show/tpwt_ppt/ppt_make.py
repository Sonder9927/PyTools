from pptx import Presentation
from pathlib import Path
from .ppt_adds import ppt_add_single_type, ppt_add_diffs


class PptMaker:
    def __init__(self, pn, fig_root: Path, remake=True) -> None:
        self.ppt_name = pn
        self.figs = fig_root
        self.margin = [1, .789]
        self.shape = [7, 8.5]
        # create a new instance for pptx
        self.prs = Presentation() if remake else Presentation(self.ppt_name)
        # make first slide
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "TPWT Results"

    def add_diffs(self, diff_dir, info_file):
        """
        add diff of all periods
        """
        self.prs = ppt_add_diffs(self.prs, self.figs / diff_dir, info_file,
                                 self.margin)

    def add_vels(self, vel_dir):
        """
        add phase vel of all periods
        """
        self.prs = ppt_add_single_type(self.prs, self.figs / vel_dir,
                                       self.margin, self.shape, "*Vel*")

    def add_CBs(self, cb_dir):
        """
        add check board of all periods
        """
        self.prs = ppt_add_single_type(self.prs, self.figs / cb_dir,
                                       self.margin, self.shape, "*CB*")

    def save(self, target=None):
        """
        save ppt
        """
        pn = target or self.ppt_name
        self.prs.save(pn)
