from pptx import util
from pathlib import Path
import json


def ppt_add_dcs(prs, figs, margin, shape):
    """
    This script can insert r*c pictures of dispersion curve
    into a new slide of pptx.
    """
    config_param = {
        "rc_n": [5, 3],
        "rc_i": [0.382, 1],
        "margin": margin,
        "shape": shape,
        "title": "Dispersion Curves",
    }
    # make batch figs for every slide
    dc_figs = sorted(figs.rglob("*.png"))
    return slide_add_batch_with_title(prs, dc_figs, config_param)


def ppt_add_probs(prs, figs, margin, shape):
    """
    This script can insert r*c pictures of
    probalCrs of all grids generated from mc
    into a new slide of pptx.
    """
    config_param = {
        "rc_n": [2, 3],
        "rc_i": [0.8, 0.5],
        "margin": margin,
        "shape": shape,
        "title": "probalCrs by mcmc",
    }
    # make batch figs for every slide
    prob_figs = sorted(figs.rglob("*/*prob*"))
    return slide_add_batch_with_title(prs, prob_figs, config_param)


def ppt_add_one_fig_per_slide(prs, figs, margin, idt):
    [left, top] = margin
    blank_slide_layout = prs.slide_layouts[6]

    for f in figs.glob(idt):
        # insert one fig into new slide
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(str(f), util.Cm(left), util.Cm(top))

    return prs


def ppt_add_diffs(prs, figs: Path, info_file, margin):
    """
    This script can insert a diff picture
    with info into a new slide of pptx.
    """
    [left, top] = margin
    blank_slide_layout = prs.slide_layouts[6]
    # read vel info json
    with open(info_file) as ff:
        info = json.load(ff)

    for f in figs.glob("*diff*"):
        # insert diff fig into new slide
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(str(f), util.Cm(left), util.Cm(top))
        # insert table
        per = f.stem.split("_")[-1]
        table_rows, table_cols = 3, 2
        table_left, table_top = util.Inches(4.5), util.Inches(4)
        table_width, table_height = util.Inches(3.5), util.Inches(0.8)
        table = slide.shapes.add_table(
            table_rows,
            table_cols,
            table_left,
            table_top,
            table_width,
            table_height,
        ).table

        table.style = "VALUE TABLE"
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"

        table.cell(1, 0).text = "mean"
        table.cell(1, 1).text = info[per]["avg_diff"]
        table.cell(2, 0).text = "std"
        table.cell(2, 1).text = info[per]["standard_deviation"]

    return prs


def ppt_add_single_type(prs, figs, margin, shape, idt, key):
    """
    This script can insert r*c pictures
    into a new slide of pptx.
    """
    config_param = {
        "rc_n": [2, 3],
        "rc_i": [0.382, 1],
        "margin": margin,
        "shape": shape,
        "title": f"{idt.replace(r'*','')} Results",
    }
    # make batch figs for every slide
    pictures = sorted(figs.glob(idt), key=key)
    return slide_add_batch_with_title(prs, pictures, config_param)


def slide_add_batch_with_title(prs, figs, config):
    nrow, ncol = config["rc_n"]  # shape
    irow, icol = config["rc_i"]  # itvel
    [left_margin, top_margin] = config["margin"]
    [width, height] = config["shape"]

    # insert title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = config["title"]

    batch = nrow * ncol
    figs_batch = [figs[i: i + batch] for i in range(0, len(figs), batch)]
    blank_slide_layout = prs.slide_layouts[6]
    for fb in figs_batch:
        # new slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # add pictures
        for i, f in enumerate(fb):
            # set margin for each fig
            left = util.Cm(left_margin + i % ncol * (width + icol))
            top = util.Cm(top_margin + i // ncol * (height + irow))

            slide.shapes.add_picture(
                str(f), left, top, util.Cm(width), util.Cm(height)
            )

    return prs
