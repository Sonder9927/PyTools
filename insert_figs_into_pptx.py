# Author : Sonder
# Created : 30 June 2022
# Version : 1.0

"""
description
    This script can insert 2*3 pictures into a new slide of pptx.
"""
from pptx import Presentation, util
import os
from icecream import ic

def plot_ts(damping, snr, smooth, width, height):

    # insert a blank template
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    for tcutoff in params["tcutoff"]:
        # find boundary fig
        fig_dir = f"{FIG_PATH}/fig_snr{snr}_tcutoff{tcutoff}_smooth{smooth}_damping{damping}"
        for f in os.listdir(fig_dir):
            if "boundary" in f:
                fig = f
                break

        # set margin for each fig
        top = util.Cm(1.382 + tcutoff // 7 * (height+.382))
        left = util.Cm(.789 + (tcutoff / 2 - 1) * (width+1.618) - tcutoff // 7 * (3 * (width+1.618))) # dui qi

        slide.shapes.add_picture(f"{fig_dir}/{fig}", left, top, util.Cm(width), util.Cm(height))

def add_figs():

    # declare global var
    global prs, params, FIG_PATH

    #directory of loopFig
    FIG_PATH = "d_num/Param/loopFig_3.65TO3.8"
    #test parameter
    params = {
        "damping": [0.2, 0.25],
        "snr": [10, 12, 15],
        "tcutoff": [4, 6, 8, 10, 12],
        "smooth": [50, 65, 80],
    }

    ic("open ppt")
    # initial a instance
    prs = Presentation(r"dp_test.pptx")

    # width and height of pictures
    width, height = 9, 8.5

    # start to insert figs
    for damping in params["damping"]:
        for snr in params["snr"]:
            for smooth in params["smooth"]:
                plot_ts(damping, snr, smooth, width, height)

    # save ppt
    prs.save("dp_test.pptx")
    ic("Good luck!")

if __name__ == "__main__":
    add_figs()
